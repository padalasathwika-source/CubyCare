import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    PRIMARY = RGBColor(0, 119, 182)      # Ocean Blue
    PRIMARY_DARK = RGBColor(2, 62, 138)   # Deep Navy
    ACCENT = RGBColor(255, 158, 0)       # Warm Gold/Orange
    TEXT_DARK = RGBColor(15, 23, 42)     # Slate Dark
    TEXT_MUTED = RGBColor(71, 85, 105)   # Muted Grey
    BG_CARD = RGBColor(241, 245, 249)    # Light Slate Card
    WHITE = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(203, 213, 225)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="CUBYCARE PLATFORM"):
        # Header banner shape
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY_DARK

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Background card
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_DARK
    bg.line.fill.background()

    # Title Box
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "CubyCare"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    p2 = tf1.add_paragraph()
    p2.text = "Complete Smart Child Healthcare, Growth & Milestone Tracking Platform"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(14)
    
    p3 = tf1.add_paragraph()
    p3.text = "An All-in-One Pediatric Companion for Modern Digital Parenting"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(224, 242, 254)
    p3.space_before = Pt(20)

    # Footer callout
    p4 = tf1.add_paragraph()
    p4.text = "Auditorium Feature Presentation | Smart Pediatric Health Ecosystem"
    p4.font.size = Pt(13)
    p4.font.color.rgb = RGBColor(186, 230, 253)
    p4.space_before = Pt(30)


    # ==================== SLIDE 2: PROBLEM & SOLUTION ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "The Problem & The CubyCare Solution")

    # Card 1: Problem
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(254, 242, 242)
    card1.line.color.rgb = RGBColor(252, 165, 165)
    
    tf = card1.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "CHALLENGES IN CHILDCARE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)
    
    problems = [
        ("Fragmented Health Records", "Paper-based immunization cards often get lost, leading to missed mandatory vaccines."),
        ("Developmental Ambiguity", "Parents struggle to track growth benchmarks (height/weight) against WHO standards."),
        ("Multi-Child Confusion", "Managing profiles, allergies, and daily medicine doses for multiple children creates friction."),
        ("Account Data Leakage", "Traditional apps mix profiles when switching parent logins or reinstalling.")
    ]
    for title, desc in problems:
        p_t = tf.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.bold = True
        p_t.font.size = Pt(13)
        p_t.font.color.rgb = TEXT_DARK
        p_t.space_before = Pt(10)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    # Card 2: Solution
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(240, 253, 244)
    card2.line.color.rgb = RGBColor(134, 239, 172)
    
    tf = card2.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "THE CUBYCARE SOLUTION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 128, 61)
    
    solutions = [
        ("Unified Pediatric Dashboard", "Single tap access to vaccines, growth logs, food diary, medicines, and vitals."),
        ("Exact Age Tracking", "Live age display computed dynamically down to exact Years, Months, and Days."),
        ("Isolated Parent Accounts", "Strict Firebase Firestore mapping ensures parent data remains completely separate."),
        ("Emergency Symptoms Engine", "CubyAlert system provides immediate red-flag condition triage and guidance.")
    ]
    for title, desc in solutions:
        p_t = tf.add_paragraph()
        p_t.text = f"✔ {title}"
        p_t.font.bold = True
        p_t.font.size = Pt(13)
        p_t.font.color.rgb = TEXT_DARK
        p_t.space_before = Pt(10)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED


    # ==================== SLIDE 3: 6 CORE PILLARS ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "The 6 Core Pillars of CubyCare")

    pillars = [
        ("1. Dynamic Child Profile", "Calculates age in Years, Months & Days. Secure avatar upload, gender & birth metrics logging."),
        ("2. Growth & Milestones", "WHO percentile growth curves (Weight/Height/Head) + age-based cognitive milestone checklists."),
        ("3. Smart Vaccination", "Complete WHO/National immunization schedule, status tracker, and automated push reminders."),
        ("4. Nutrition & Food Diary", "Daily meal tracker, allergy management, dietary preferences, and nutritional history."),
        ("5. Medicine & Symptoms", "Active dosage schedules, medicine flowchart, treatment history, and CubyAlert triage."),
        ("6. Play, Naps & Guide", "Age-appropriate fun activities, sleep duration benchmarks, and expert parenting advice.")
    ]

    for idx, (p_title, p_desc) in enumerate(pillars):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.6 + row * 2.6)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = CARD_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = p_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)


    # ==================== SLIDE 4: SMART HEALTH & EMERGENCY ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Smart Health, Medicines & Emergency Care")

    features = [
        ("Medicine Tracker & Dosage Schedule", [
            "Track active vs completed medication courses.",
            "Record exact dosage, frequency, and prescribing doctor.",
            "Interactive Medicine Flowchart for treatment history."
        ]),
        ("CubyAlert Emergency Symptoms Engine", [
            "Instant red-flag triage for High Fever, Dehydration, Breathing issues.",
            "Clear step-by-step first-aid and emergency call triggers.",
            "Peace of mind during critical pediatric moments."
        ]),
        ("Vitals & Doctor Appointment Manager", [
            "Log body temperature, heart rate, oxygen level (SpO2).",
            "Store pediatrician appointment notes and follow-up dates.",
            "Comprehensive health summary ready for hospital visits."
        ])
    ]

    for idx, (title, points) in enumerate(features):
        x = Inches(0.8 + idx * 3.9)
        y = Inches(1.6)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PRIMARY
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        
        for pt in points:
            p_bullet = tf.add_paragraph()
            p_bullet.text = f"• {pt}"
            p_bullet.font.size = Pt(11)
            p_bullet.font.color.rgb = TEXT_DARK
            p_bullet.space_before = Pt(12)


    # ==================== SLIDE 5: SECURITY & ACCOUNT ISOLATION ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Architecture: Privacy, Account Isolation & Sync")

    # Left Card - Security & Isolation
    card_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = BG_CARD
    card_l.line.color.rgb = PRIMARY
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = Inches(0.3)
    
    p = tf_l.paragraphs[0]
    p.text = "STRICT ACCOUNT ISOLATION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    
    iso_points = [
        ("Scoped Firestore Collections", "Every child profile is mapped to the logged-in parent's unique Auth UID."),
        ("Zero Cross-Account Leakage", "Logging in with a different account or reinstalling the app cleanly isolated child profiles."),
        ("Parent Account Removal", "Full capability to remove parent accounts and delete associated data securely."),
        ("Guest Mode Handling", "Temporary local sandbox that does not contaminate registered cloud parent profiles.")
    ]
    for t, d in iso_points:
        p_t = tf_l.add_paragraph()
        p_t.text = f"🔒 {t}"
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = TEXT_DARK
        p_t.space_before = Pt(10)
        p_d = tf_l.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    # Right Card - Data Sync & Offline First
    card_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = BG_CARD
    card_r.line.color.rgb = PRIMARY
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "CLOUD SYNC & OFFLINE RESILIENCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    
    sync_points = [
        ("Real-time Firebase Firestore Sync", "Instant synchronization across mother & father devices when connected."),
        ("Offline Local Cache", "Local persistence ensures vaccination & growth records are accessible without internet."),
        ("Android WorkManager Integration", "Background vaccine reminder notifications scheduled reliably 24/7."),
        ("Encrypted Security", "Industry-standard Firebase Authentication protecting health data.")
    ]
    for t, d in sync_points:
        p_t = tf_r.add_paragraph()
        p_t.text = f"☁️ {t}"
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = TEXT_DARK
        p_t.space_before = Pt(10)
        p_d = tf_r.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED


    # ==================== SLIDE 6: UX EXCELLENCE & THEMES ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "UI/UX Design, Theme Engine & Smooth Performance")

    ux_items = [
        ("Silky Material Crossfade", "Eliminated frame drops with FastOutSlowInEasing (200ms) transitions for fluid navigation."),
        ("Custom Color Palette Picker", "Multiple aesthetic themes: Ocean Blue, Emerald Green, Soft Rose, and Sunset Warmth."),
        ("High-Legibility Dark Mode", "True dark palette with #0F172A background and high-contrast #F8FAFC text tokens."),
        ("Enhanced Settings Controls", "Interactive ON/OFF visual pill toggles for notifications, reminders, and sound."),
        ("Vivid Ambient Artwork", "Soft, non-distracting background artwork tailored to nursery and mother-baby care.")
    ]

    for idx, (title, desc) in enumerate(ux_items):
        y = Inches(1.6 + idx * 1.05)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_DARK


    # ==================== SLIDE 7: ROADMAP & CONCLUSION ====================
    slide7 = prs.slides.add_slide(blank_layout)
    
    # Background card
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = PRIMARY_DARK
    bg7.line.fill.background()

    tb7 = slide7.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(6.0))
    tf7 = tb7.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "FUTURE ROADMAP & IMPACT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    p_t = tf7.add_paragraph()
    p_t.text = "Transforming Digital Child Healthcare"
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.space_before = Pt(6)

    roadmap_points = [
        ("AI Pediatric Assistant (Gemini API)", "Smart conversational AI to answer parenting queries and analyze symptoms."),
        ("Pediatrician PDF Export", "One-click export of complete growth & vaccine records for hospital visits."),
        ("Smart Wearable Sync", "Direct Bluetooth connectivity with smart thermometers & pulse oximeters."),
        ("Community & Expert Forums", "Verified pediatrician Q&A space for milestone guidance.")
    ]

    for title, desc in roadmap_points:
        p_r = tf7.add_paragraph()
        p_r.text = f"🚀 {title} — {desc}"
        p_r.font.size = Pt(14)
        p_r.font.color.rgb = RGBColor(224, 242, 254)
        p_r.space_before = Pt(16)

    p_close = tf7.add_paragraph()
    p_close.text = "CubyCare — Nurturing Every Milestone with Love, Precision & Peace of Mind."
    p_close.font.size = Pt(18)
    p_close.font.bold = True
    p_close.font.color.rgb = ACCENT
    p_close.space_before = Pt(36)

    # Save presentation
    output_path = "CubyCare_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
